"""Generate Suistody hackathon presentation PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Theme colors --
VOID = RGBColor(0x06, 0x0A, 0x13)
DEEP = RGBColor(0x0A, 0x0E, 0x1A)
ACCENT = RGBColor(0x00, 0xD4, 0xFF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
LIGHT_GRAY = RGBColor(0xD1, 0xD5, 0xDB)
DARK_BORDER = RGBColor(0x1F, 0x29, 0x37)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, font_name="Segoe UI", space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_accent_line(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.8), top, Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, VOID)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
         "SUISTODY", 72, ACCENT, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3.0), Inches(11), Inches(1),
         '"Don\'t give your AI agent the keys. Give it a budget."',
         28, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
         "Policy-Based AI Agent Custody on Sui",
         24, WHITE, True, PP_ALIGN.CENTER)

# Hackathon badge
add_shape_rect(slide, Inches(4.2), Inches(5.5), Inches(5), Inches(0.7),
               DEEP, ACCENT)
add_text(slide, Inches(4.2), Inches(5.55), Inches(5), Inches(0.6),
         "Sui Vibe Hackathon 2026  |  Cetus + Stablelayer + Sui Track",
         16, ACCENT, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "github.com/ARZER-TW/agent-vault  |  agent-vault-dusky.vercel.app",
         14, GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "THE PROBLEM", 40, WHITE, True)

add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1),
         "AI Agents increasingly need to transact autonomously --\ncalling APIs, purchasing cloud resources, executing DeFi trades.",
         22, LIGHT_GRAY, False)

# Problem cards
problems = [
    ("Full Key Access", "Give agents private keys\n= catastrophic risk", RED),
    ("Human Approval", "Require approval every TX\n= defeats autonomy", AMBER),
    ("EVM Approve", "Only amount cap, no action /\ncooldown / expiry control", GRAY),
]
for i, (title, desc, color) in enumerate(problems):
    x = Inches(0.8 + i * 4)
    add_shape_rect(slide, x, Inches(3.5), Inches(3.5), Inches(2.5), DEEP, color)
    add_text(slide, x + Inches(0.3), Inches(3.7), Inches(2.9), Inches(0.6),
             title, 22, color, True)
    add_text(slide, x + Inches(0.3), Inches(4.4), Inches(2.9), Inches(1.4),
             desc, 18, LIGHT_GRAY)


# ============================================================
# SLIDE 3: The Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "THE SOLUTION", 40, WHITE, True)

add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.8),
         "Create a Vault with multi-dimensional policy. Give your agent a budget, not your keys.",
         22, LIGHT_GRAY)

# 5 Policy dimensions
policies = [
    ("Max Budget", "Total spending cap", ACCENT),
    ("Max Per TX", "Per-transaction limit", ACCENT),
    ("Allowed Actions", "Whitelist operations", GREEN),
    ("Cooldown", "Min time between TXs", AMBER),
    ("Expiration", "Auto-revoke deadline", RED),
]
for i, (title, desc, color) in enumerate(policies):
    x = Inches(0.5 + i * 2.5)
    add_shape_rect(slide, x, Inches(3.3), Inches(2.2), Inches(1.8), DEEP, color)
    add_text(slide, x + Inches(0.2), Inches(3.5), Inches(1.8), Inches(0.5),
             title, 18, color, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(4.1), Inches(1.8), Inches(0.8),
             desc, 15, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(1.2),
         "AgentCap = transferable NFT permission token\n"
         "Every withdrawal validated against ALL 5 dimensions atomically on-chain\n"
         "Owner can revoke AgentCap instantly at any time",
         18, GRAY)


# ============================================================
# SLIDE 4: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "ARCHITECTURE", 40, WHITE, True)

# 4 layers
layers = [
    ("Frontend", "Next.js 14 + React 18\nVault Noir Design System\nzkLogin (Google OAuth)", ACCENT, Inches(0.5)),
    ("AI Agent Runtime", "Multi-LLM (GPT-4o / Gemini / Claude)\nZod Intent Parser\n7-Step Pipeline", PURPLE, Inches(3.5)),
    ("Policy Engine", "Off-Chain Pre-check (6 rules)\nOn-Chain Enforcement (9 checks)\nDual-Layer Security", AMBER, Inches(6.5)),
    ("Sui Blockchain", "Move Smart Contract\nCetus Aggregator SDK\nStablelayer SDK", GREEN, Inches(9.5)),
]
for title, desc, color, x in layers:
    add_shape_rect(slide, x, Inches(2.2), Inches(3), Inches(4), DEEP, color)
    add_text(slide, x + Inches(0.25), Inches(2.4), Inches(2.5), Inches(0.5),
             title, 20, color, True, PP_ALIGN.CENTER)
    # separator line
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x + Inches(0.5), Inches(3.0), Inches(2), Pt(1.5))
    sep.fill.solid()
    sep.fill.fore_color.rgb = color
    sep.line.fill.background()
    add_text(slide, x + Inches(0.25), Inches(3.2), Inches(2.5), Inches(2.5),
             desc, 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: Agent Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "AI AGENT 7-STEP PIPELINE", 40, WHITE, True)

steps = [
    ("1", "Fetch\nMarket Data", ACCENT),
    ("2", "Query\nLLM", PURPLE),
    ("3", "Parse\nIntent", LIGHT_GRAY),
    ("4", "Policy\nPre-Check", AMBER),
    ("5", "Build\nPTB", GREEN),
    ("6", "Sponsored\nExecution", ACCENT),
    ("7", "Log\nResult", GRAY),
]
for i, (num, desc, color) in enumerate(steps):
    x = Inches(0.3 + i * 1.8)
    # Circle number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.45), Inches(2.5),
                                     Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf_c = circle.text_frame
    tf_c.paragraphs[0].text = num
    tf_c.paragraphs[0].font.size = Pt(28)
    tf_c.paragraphs[0].font.color.rgb = VOID
    tf_c.paragraphs[0].font.bold = True
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_c.word_wrap = False

    add_text(slide, x, Inches(3.6), Inches(1.7), Inches(1),
             desc, 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # Arrow between steps
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        x + Inches(1.35), Inches(2.7),
                                        Inches(0.4), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_BORDER
        arrow.line.fill.background()

# Bottom note
add_text(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5),
         "Natural Language Strategy: Tell AI how to trade in plain English\n"
         "4 Presets: Conservative DCA | Take Profit | Aggressive Trading | Minimal Risk\n"
         "Auto-Run Mode: 30s / 45s / 60s / 120s intervals with live activity log",
         17, GRAY)


# ============================================================
# SLIDE 6: Security - Dual Layer
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "DUAL-LAYER SECURITY", 40, WHITE, True)

add_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
         "Policy enforced TWICE -- off-chain (save gas) + on-chain (guarantee correctness)",
         20, LIGHT_GRAY)

# Off-chain box
add_shape_rect(slide, Inches(0.5), Inches(2.6), Inches(5.8), Inches(4.2), DEEP, AMBER)
add_text(slide, Inches(0.8), Inches(2.8), Inches(5.2), Inches(0.5),
         "OFF-CHAIN PRE-CHECK (policy-checker.ts)", 18, AMBER, True)
off_checks = [
    "1. Zero Amount?",
    "2. Expired?",
    "3. Cooldown Active?",
    "4. Exceeds Per-TX Limit?",
    "5. Exceeds Total Budget?",
    "6. Action Whitelisted?",
    "7. Sufficient Balance?",
]
tf = add_text(slide, Inches(1), Inches(3.4), Inches(5), Inches(3),
              off_checks[0], 16, LIGHT_GRAY, False, font_name="Consolas")
for check in off_checks[1:]:
    add_para(tf, check, 16, LIGHT_GRAY, font_name="Consolas")

# On-chain box
add_shape_rect(slide, Inches(7), Inches(2.6), Inches(5.8), Inches(4.2), DEEP, GREEN)
add_text(slide, Inches(7.3), Inches(2.8), Inches(5.2), Inches(0.5),
         "ON-CHAIN ENFORCEMENT (agent_vault.move)", 18, GREEN, True)
on_checks = [
    "1. assert amount > 0",
    "2. assert cap.vault_id == vault",
    "3. assert cap in authorized_caps",
    "4. assert now < expires_at",
    "5. assert cooldown elapsed",
    "6. assert amount <= max_per_tx",
    "7. assert amount <= budget - spent",
    "8. assert action in allowed_actions",
    "9. assert balance >= amount",
]
tf2 = add_text(slide, Inches(7.5), Inches(3.4), Inches(5), Inches(3),
               on_checks[0], 16, LIGHT_GRAY, False, font_name="Consolas")
for check in on_checks[1:]:
    add_para(tf2, check, 15, LIGHT_GRAY, font_name="Consolas", space_before=Pt(3))


# ============================================================
# SLIDE 7: Why Sui?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "WHY SUI? (Can't Be Built on EVM)", 40, WHITE, True)

features = [
    ("Object Capabilities", "AgentCap = 5D policy object\nvs EVM approve() = amount only", ACCENT),
    ("PTB", "withdraw + swap + transfer\nin ONE atomic TX, ONE gas fee", GREEN),
    ("zkLogin", "Google login = Sui address\nNo MetaMask, no seed phrase", PURPLE),
    ("Sponsored TX", "Zero gas for users AND agents\nNative protocol support", AMBER),
    ("Move Type Safety", "AgentCap can't be copied\nCompiler-enforced, not runtime", RED),
]
for i, (title, desc, color) in enumerate(features):
    y = Inches(2.0 + i * 1.05)
    # Color dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.1),
                                  Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(slide, Inches(1.3), y, Inches(3.5), Inches(0.4),
             title, 22, color, True)
    add_text(slide, Inches(5), y, Inches(7.5), Inches(0.9),
             desc, 17, LIGHT_GRAY)


# ============================================================
# SLIDE 8: Cetus & Stablelayer Integration
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "CETUS & STABLELAYER INTEGRATION", 40, WHITE, True)

# Cetus box
add_shape_rect(slide, Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.5), DEEP, ACCENT)
add_text(slide, Inches(0.8), Inches(2.4), Inches(5.2), Inches(0.5),
         "Cetus Aggregator SDK", 24, ACCENT, True)
cetus_items = [
    "@cetusprotocol/aggregator-sdk v1.4.4",
    "Cross 25+ DEX route aggregation",
    "findRouters() for optimal swap path",
    "routerSwap() for on-chain execution",
    "1% default slippage tolerance",
    "Auto-fallback to simple withdraw",
]
tf = add_text(slide, Inches(1), Inches(3.2), Inches(5), Inches(3),
              cetus_items[0], 16, LIGHT_GRAY, False, font_name="Consolas")
for item in cetus_items[1:]:
    add_para(tf, item, 16, LIGHT_GRAY, font_name="Consolas")

# Stablelayer box
add_shape_rect(slide, Inches(7), Inches(2.2), Inches(5.8), Inches(4.5), DEEP, AMBER)
add_text(slide, Inches(7.3), Inches(2.4), Inches(5.2), Inches(0.5),
         "Stablelayer SDK", 24, AMBER, True)
stable_items = [
    "stable-layer-sdk v2.0.0",
    "By Bucket Protocol",
    "buildMintTx: Mint LakeUSDC",
    "buildBurnTx: Burn LakeUSDC",
    "buildClaimTx: Claim rewards",
    "Mainnet-only (code ready)",
]
tf2 = add_text(slide, Inches(7.5), Inches(3.2), Inches(5), Inches(3),
               stable_items[0], 16, LIGHT_GRAY, False, font_name="Consolas")
for item in stable_items[1:]:
    add_para(tf2, item, 16, LIGHT_GRAY, font_name="Consolas")


# ============================================================
# SLIDE 9: Beyond DeFi
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "BEYOND DeFi -- THE BIGGER PICTURE", 40, WHITE, True)

add_text(slide, Inches(0.8), Inches(1.9), Inches(11), Inches(0.7),
         "Suistody is not a DeFi tool. It's a universal permission layer for autonomous AI agents.",
         22, LIGHT_GRAY)

# Vision statement
add_shape_rect(slide, Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.9), DEEP, ACCENT)
add_text(slide, Inches(0.8), Inches(2.9), Inches(11.7), Inches(0.7),
         '"Any AI agent that needs to spend money autonomously, but shouldn\'t have unlimited access."',
         20, ACCENT, True, PP_ALIGN.CENTER)

# Use cases grid - 2 rows x 3 columns
use_cases = [
    ("AI Autonomous Payments", "Agents buy API credits,\ncloud resources, subscriptions\nwith daily/monthly caps", ACCENT),
    ("DAO Treasury", "AI manages DAO funds,\nexecutes approved proposals\nwithin voted budgets", GREEN),
    ("Gaming", "AI controls in-game assets,\nbuys/sells with spending\nlimits per session", PURPLE),
    ("NFT Trading", "AI auto-trades NFTs\nby strategy, constrained by\nper-TX and total budget", AMBER),
    ("Infrastructure", "AI pays for decentralized\ncompute, storage, bandwidth\nwith cooldown controls", LIGHT_GRAY),
    ("Social & Tipping", "AI rewards creators,\ntips content, donates --\nall within daily caps", RED),
]
for i, (title, desc, color) in enumerate(use_cases):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(4.0 + row * 1.7)
    add_shape_rect(slide, x, y, Inches(3.8), Inches(1.5), DEEP, color)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.4), Inches(0.4),
             title, 16, color, True)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.4), Inches(0.9),
             desc, 14, LIGHT_GRAY)


# ============================================================
# SLIDE 10: Guardrail Stress Test
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "GUARDRAIL STRESS TEST", 40, WHITE, True)

add_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
         "5 adversarial scenarios -- ALL must be BLOCKED for a correctly configured vault",
         20, LIGHT_GRAY)

tests = [
    ("1", "Budget Overflow", "Exceed remaining budget", RED),
    ("2", "Per-TX Breach", "Exceed per-transaction limit", RED),
    ("3", "Cooldown Bypass", "Trade during cooldown period", RED),
    ("4", "Unauthorized Agent", "Use non-authorized AgentCap", RED),
    ("5", "Expired Policy", "Trade after policy expiry", RED),
]
for i, (num, title, desc, color) in enumerate(tests):
    y = Inches(2.7 + i * 0.95)
    add_shape_rect(slide, Inches(1), y, Inches(11), Inches(0.8), DEEP, DARK_BORDER)
    # BLOCKED badge
    add_shape_rect(slide, Inches(1.3), y + Inches(0.15), Inches(1.5), Inches(0.5),
                   RGBColor(0x30, 0x10, 0x10), RED)
    add_text(slide, Inches(1.3), y + Inches(0.15), Inches(1.5), Inches(0.5),
             "BLOCKED", 14, RED, True, PP_ALIGN.CENTER, "Consolas")
    add_text(slide, Inches(3.2), y + Inches(0.1), Inches(3), Inches(0.6),
             f"{num}. {title}", 19, WHITE, True)
    add_text(slide, Inches(7), y + Inches(0.15), Inches(4.5), Inches(0.5),
             desc, 16, GRAY)


# ============================================================
# SLIDE 10: Tech Stack
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "TECH STACK", 40, WHITE, True)

stack = [
    ("Frontend", "Next.js 14 + TypeScript + Tailwind CSS", ACCENT),
    ("State", "Zustand 5 + React Query 5", ACCENT),
    ("Sui SDK", "@mysten/sui v1.44.0", GREEN),
    ("DeFi", "Cetus Aggregator v1.4.4 + Stablelayer v2.0.0", AMBER),
    ("AI", "GPT-4o | Gemini 2.0 Flash | Claude Sonnet (auto-detect)", PURPLE),
    ("Auth", "zkLogin (Google OAuth + Enoki ZK Prover)", PURPLE),
    ("Contracts", "Sui Move (edition 2024.beta)", GREEN),
    ("Validation", "Zod v3.24 (all LLM responses validated)", LIGHT_GRAY),
    ("Testing", "Vitest (78 tests) + sui move test (15 tests)", LIGHT_GRAY),
]
for i, (label, desc, color) in enumerate(stack):
    y = Inches(2.0 + i * 0.58)
    add_text(slide, Inches(1), y, Inches(3), Inches(0.5),
             label, 18, color, True)
    add_text(slide, Inches(4), y, Inches(8.5), Inches(0.5),
             desc, 17, LIGHT_GRAY)


# ============================================================
# SLIDE 11: Test Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "TEST RESULTS", 40, WHITE, True)

# Test stats
stats = [
    ("78/78", "TypeScript\nUnit Tests", ACCENT),
    ("15/15", "Move Contract\nTests", GREEN),
    ("5/5", "Guardrail\nStress Tests", AMBER),
    ("93/93", "Total Tests\nAll Passing", WHITE),
]
for i, (num, label, color) in enumerate(stats):
    x = Inches(0.8 + i * 3.1)
    add_shape_rect(slide, x, Inches(2.5), Inches(2.6), Inches(2.5), DEEP, color)
    add_text(slide, x, Inches(2.8), Inches(2.6), Inches(1),
             num, 48, color, True, PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.9), Inches(2.6), Inches(0.8),
             label, 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Coverage list
add_text(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
         "Test Coverage: intent-parser (20) + policy-checker (14) + ptb-builder (13) + "
         "ptb-agent (6) + service (14) + constants (11) + Move contract (15)",
         16, GRAY)


# ============================================================
# SLIDE 12: AI Disclosure
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)
add_accent_line(slide, Inches(0.8))
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "AI TOOL DISCLOSURE", 40, WHITE, True)

add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.5),
         "As required by hackathon rules, full transparency on AI tools used:",
         20, LIGHT_GRAY)

disclosures = [
    ("Tool", "Claude Code (CLI)", ACCENT),
    ("Model", "Claude Opus 4.6 (claude-opus-4-6)", ACCENT),
    ("Usage", "Architecture design, code generation,\ndebugging, test writing, documentation", LIGHT_GRAY),
    ("Key Prompts", "Implementation planning, TDD workflow,\nMove contract design, SDK integration", LIGHT_GRAY),
    ("Note", "All code reviewed and tested by developer.\n93 tests passing. Full open source.", GREEN),
]
for i, (label, desc, color) in enumerate(disclosures):
    y = Inches(3.0 + i * 0.85)
    add_text(slide, Inches(1), y, Inches(3), Inches(0.7),
             label, 20, AMBER, True)
    add_text(slide, Inches(4), y, Inches(8.5), Inches(0.7),
             desc, 18, color)


# ============================================================
# SLIDE 13: Live Demo + Links
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, VOID)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
         "LIVE DEMO", 56, ACCENT, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
         "agent-vault-dusky.vercel.app",
         32, WHITE, True, PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(3), Inches(3.8), Inches(7.3), Inches(2.5), DEEP, ACCENT)
tf = add_text(slide, Inches(3.5), Inches(4.0), Inches(6.3), Inches(2),
              "1. Sign in with Google (zkLogin)", 20, LIGHT_GRAY)
add_para(tf, "2. Create a Vault with policy", 20, LIGHT_GRAY)
add_para(tf, "3. Run Agent cycle (AI or Demo mode)", 20, LIGHT_GRAY)
add_para(tf, "4. Run Guardrail Stress Test", 20, LIGHT_GRAY)
add_para(tf, "5. View On-Chain Audit Trail on SuiScan", 20, LIGHT_GRAY)

# Footer
add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "GitHub: github.com/ARZER-TW/agent-vault    |    Sui Vibe Hackathon 2026",
         16, GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# Save
# ============================================================
output_path = "/home/james/hackathons/agent-vault/Suistody_Presentation.pptx"
prs.save(output_path)
print(f"[OK] Saved to {output_path}")
